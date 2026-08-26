"""SARIMA.ipynb ciktilarindan SUNUM1 PowerPoint dosyasini uretir."""

from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SUNUM_DIR = Path(__file__).resolve().parent
IMAGE_DIR = SUNUM_DIR / "gorseller"

DARK = RGBColor(0x1F, 0x29, 0x37)
ANT = RGBColor(0xE8, 0x59, 0x0C)     # Antalya turuncu
MUG = RGBColor(0x19, 0x71, 0xC2)     # Mugla mavi
WARN = RGBColor(0x9C, 0x27, 0x11)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF1, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def img(name):
    path = IMAGE_DIR / name
    return path if path.exists() else None


def add_band(slide, title, subtitle=None, color=DARK):
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.02))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    band.shadow.inherit = False
    frame = band.text_frame
    frame.margin_left = Inches(0.5)
    frame.margin_top = Inches(0.1)
    frame.text = title
    run = frame.paragraphs[0].runs[0]
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = WHITE
    frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    if subtitle:
        para = frame.add_paragraph()
        para.text = subtitle
        para.runs[0].font.size = Pt(12.5)
        para.runs[0].font.color.rgb = RGBColor(0xCE, 0xD4, 0xDA)


def add_text(slide, left, top, width, height, lines, size=15, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        indented = line.startswith("    ")
        text = line.strip()
        if text.startswith("- "):
            para.text = "   • " + text[2:]
        else:
            para.text = text
        para.space_after = Pt(6)
        for run in para.runs:
            run.font.size = Pt(size - 2 if indented else size)
            run.font.color.rgb = MUTED if indented else color
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h):
    with Image.open(path) as image:
        width_px, height_px = image.size
    scale = min(max_w / width_px, max_h / height_px)
    width, height = Emu(int(width_px * scale)), Emu(int(height_px * scale))
    return slide.shapes.add_picture(
        str(path), Emu(int(left + (max_w - width) / 2)), Emu(int(top + (max_h - height) / 2)),
        width, height,
    )


def add_table(slide, frame_data, left, top, width, height, col_widths=None, size=11):
    shape = slide.shapes.add_table(frame_data.shape[0] + 1, frame_data.shape[1], left, top, width, height)
    table = shape.table
    if col_widths:
        for index, ratio in enumerate(col_widths):
            table.columns[index].width = Emu(int(width * ratio))
    for col_index, column_name in enumerate(frame_data.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = WHITE
    for row_index, (_, record) in enumerate(frame_data.iterrows(), start=1):
        for col_index, value in enumerate(record):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else LIGHT
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(size)
            run.font.color.rgb = DARK
    return table


def add_kpi(slide, left, top, width, height, label, value, sub, color):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = color
    card.line.width = Pt(2)
    card.shadow.inherit = False
    frame = card.text_frame
    frame.margin_left = Inches(0.22)
    frame.margin_top = Inches(0.16)
    frame.word_wrap = True
    frame.text = label
    frame.paragraphs[0].runs[0].font.size = Pt(11.5)
    frame.paragraphs[0].runs[0].font.color.rgb = MUTED
    value_para = frame.add_paragraph()
    value_para.text = value
    value_para.runs[0].font.size = Pt(28)
    value_para.runs[0].font.bold = True
    value_para.runs[0].font.color.rgb = color
    sub_para = frame.add_paragraph()
    sub_para.text = sub
    sub_para.runs[0].font.size = Pt(10.5)
    sub_para.runs[0].font.color.rgb = DARK


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_bullets(title, subtitle, lines, note, size=15, color=DARK):
    slide = prs.slides.add_slide(BLANK)
    add_band(slide, title, subtitle, color=color)
    add_text(slide, Inches(0.65), Inches(1.35), Inches(12.1), Inches(5.7), lines, size=size)
    notes(slide, note)
    return slide


def slide_image(title, subtitle, image_name, side_lines, note, color=DARK):
    slide = prs.slides.add_slide(BLANK)
    add_band(slide, title, subtitle, color=color)
    path = img(image_name)
    if path:
        add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(8.55), Inches(5.95))
    add_text(slide, Inches(9.15), Inches(1.4), Inches(3.95), Inches(5.7), side_lines, size=12.5)
    notes(slide, note)
    return slide


def slide_full_image(title, subtitle, image_name, note, color=DARK):
    slide = prs.slides.add_slide(BLANK)
    add_band(slide, title, subtitle, color=color)
    path = img(image_name)
    if path:
        add_picture_fit(slide, path, Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.05))
    notes(slide, note)
    return slide


# ======================================================================================
# 1. Kapak
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK
bg.line.fill.background(); bg.shadow.inherit = False
stripe = slide.shapes.add_shape(1, 0, Inches(3.0), SLIDE_W, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ANT
stripe.line.fill.background(); stripe.shadow.inherit = False

box = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.2))
frame = box.text_frame
frame.text = "Antalya & Muğla Konaklama Talep Analizi"
frame.paragraphs[0].runs[0].font.size = Pt(40)
frame.paragraphs[0].runs[0].font.bold = True
frame.paragraphs[0].runs[0].font.color.rgb = WHITE

box = slide.shapes.add_textbox(Inches(1.0), Inches(3.25), Inches(11.3), Inches(2.4))
frame = box.text_frame
frame.word_wrap = True
for index, text in enumerate([
    "Şehir · İlçe · Otel kırılımında inhouse (occupied nights) tahmini",
    "SARIMAX ve Prophet karşılaştırması + takvim (bayram) düzeltmeleri",
    "Veri: 2023-2026 Temmuz  |  Test: 2026 Ocak-Temmuz  |  Kaynak: notebooks/SARIMA.ipynb",
    "26 Ağustos 2026",
]):
    para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
    para.text = text
    para.space_after = Pt(9)
    para.runs[0].font.size = Pt(17 if index == 0 else 13.5)
    para.runs[0].font.color.rgb = WHITE if index == 0 else RGBColor(0xAD, 0xB5, 0xBD)
notes(slide, "Amaç: iki ana destinasyonda aylık talebi şehir, ilçe ve otel seviyesinde tahmin "
             "etmek; 2026 gerçekleşmesiyle modeli sınamak ve iki şehri karşılaştırmak.")

# ======================================================================================
# 2. Yönetici özeti
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
add_band(slide, "Yönetici Özeti", "2026 Ocak-Temmuz backtest, SARIMA.ipynb sonuçları")
card_w, card_h, gap = Inches(3.0), Inches(1.7), Inches(0.22)
for index, (label, value, sub, color) in enumerate([
    ("Antalya · şehir WMAPE", "%17,6", "Takvim düzeltmeli SARIMAX", ANT),
    ("Antalya · düzeltmesiz", "%37,0", "SARIMAX baseline", ANT),
    ("Muğla · şehir WMAPE", "%19,5", "Takvim düzeltmeli SARIMAX", MUG),
    ("Muğla · düzeltmesiz", "%28,2", "SARIMAX baseline", MUG),
]):
    add_kpi(slide, Inches(0.55) + index * (card_w + gap), Inches(1.3), card_w, card_h,
            label, value, sub, color)
add_text(slide, Inches(0.65), Inches(3.35), Inches(12.1), Inches(3.8), [
    "Her iki şehirde de takvim (bayram) düzeltmesi hatayı yarı yarıya azaltıyor.",
    "- Antalya: %37,0 → %17,6 WMAPE   |   Muğla: %28,2 → %19,5 WMAPE",
    "İlçe seviyesinde en iyi birimler: Antalya Kemer %22 · Muğla Datça %15, Fethiye %16.",
    "Otel seviyesinde en iyi birimler: Antalya Sueno Deluxe Belek %32 · Muğla Divan Bodrum %19.",
    "Model seçimi birim bazında yapıldı: bazı ilçe/otelde SARIMAX, bazısında Prophet kazandı.",
    "İKİ UYARI (18. ve 30. slaytlarda detay):",
    "    1) Takvim düzeltmesi 2026 gerçeğinden hesaplandığı için %17,6 kör tahmin başarısı değildir.",
    "    2) SARIMAX uzun ufukta negatif değer üretiyor (2027 Haziran: -48.037) → 2027 için Prophet'e geçildi.",
    "Karar önerisi: şehir ve yüksek hacimli ilçe/otellerde model kullanılabilir; Marmaris ve",
    "    düşük hacimli birimlerde (Aksu, Döşemealtı, Milas) model çıktısı kullanılmamalı.",
], size=14.5)
notes(slide, "Açılış mesajı: model çalışıyor ve takvim düzeltmesi ciddi kazanç sağlıyor. "
             "Ama iki metodolojik uyarıyı baştan söyle ki güven kazan.")

# ======================================================================================
# 3. Veri ve kapsam
# ======================================================================================
slide_bullets(
    "Veri ve Kapsam",
    "complete_cleaned_data_full.csv · voucher bazlı rezervasyon kayıtları",
    [
        "Hedef değişken: inhouse = occupied nights (konaklanan gece sayısı).",
        "- Her rezervasyon check-in ile check-out arasındaki günlere açılır ve gün gün sayılır.",
        "- Böylece 'satış tarihi' değil, 'konaklama tarihi' bazlı gerçek doluluk elde edilir.",
        "- checkout ≤ checkin olan ve tarihi eksik kayıtlar elenir.",
        "",
        "Kapsam ve kırılımlar",
        "- Antalya: 2023-01 … 2026-07 aylık seri; 12 aktif ilçe; Top 20 otel",
        "- Muğla: 2022-01 … 2026-07 aylık seri; 10 aktif ilçe; Top 30 → filtreli 12 otel",
        "- Antalya ilçe eşlemesi: hotel_region_name → resmi ilçe (55 bölge, elle hazırlanmış sözlük)",
        "- Muğla ilçe eşlemesi: anahtar kelime tabanlı (Bodrum, Marmaris, Fethiye, Datça, Ortaca…)",
        "",
        "Hacim filtreleri (yüzdesel metriklerin patlamasını engellemek için)",
        "- Antalya otel: 2023, 2024 ve 2025 Ocak-Temmuz inhouse ≥ 200 gece",
        "- Muğla otel: 2022-2025 her sezon (Mart-Ekim) inhouse ≥ 100 gece",
        "- Düşük hacimli ilçeler görsellerden çıkarıldı: Aksu, Döşemealtı, Finike, Kepez, Konyaaltı",
    ],
    "Filtreler bilinçli: günde 1-2 gecelik birimlerde tek rezervasyon sapması %100 hata üretir, "
    "bu da tabloyu yanıltır.",
    size=14,
)

# ======================================================================================
# 4. Metodoloji
# ======================================================================================
slide_bullets(
    "Metodoloji",
    "İki aday model + açık takvim düzeltmesi + birim bazında otomatik seçim",
    [
        "Aday modeller",
        "- SARIMAX: (1,2,1)(1,0,0,12) şehir seviyesinde; (1,1,1)(1,0,0,12) ilçe/otel seviyesinde",
        "- Prophet: çarpımsal mevsimsellik, yıllık mevsimsellik açık, changepoint_prior 0.01-0.05",
        "- İlçe 2026 backtestinde ayrıca Ridge (ay etkileri + yıl trendi) kullanıldı",
        "",
        "Takvim düzeltmeleri — bayram kaymasını yakalamak için dört parça",
        "- Mayıs uplift (+): Kurban Bayramı'nın Mayıs'a kayması",
        "- Haziran transfer (×): talebin Mayıs'a çekilmesi sonucu Haziran düşüşü",
        "- Mart uplift (+) ve Nisan transfer (×): Ramazan Bayramı kayması",
        "",
        "Model seçimi",
        "- Her ilçe ve otel için tüm adaylar denenir, en düşük WMAPE'li model seçilir (MAE eşitlik bozucu).",
        "- Sonuç: bazı birimlerde SARIMAX, bazılarında Prophet kazandı — tek model dayatılmadı.",
        "",
        "Değerlendirme dönemleri",
        "- Şehir/otel: eğitim 2023-2025 (Muğla 2022-2025), test 2026 Ocak-Temmuz",
        "- İlçe: eğitim 2023-2024, test 2025 tüm aylar; ayrıca 2023-2025 Mar-Tem → 2026 Mar-Tem",
    ],
    "Metodolojinin güçlü yanı: tek model dayatmıyoruz, her birim için en iyisini ölçüyoruz. "
    "Takvim düzeltmesi de gizli bir regresör değil, açık ve denetlenebilir bir katsayı.",
    size=14,
)

# ======================================================================================
# 5. Metrik rehberi
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
add_band(slide, "Metrikler Nasıl Okunmalı?", "Hangi karar hangi metrikle verilir")
metric_table = pd.DataFrame([
    ["WMAPE", "Σ|hata| / Σgerçek", "Aylık isabet. Ay ay planlama ve fiyatlama için."],
    ["Toplam hata %", "|Σtahmin − Σgerçek| / Σgerçek", "Dönem toplamı. Kapasite ve kontenjan için."],
    ["MAPE", "Ort(|hata| / gerçek)", "Küçük aylarda patlar; düşük hacimde yanıltıcı."],
    ["R²", "Açıklanan varyans", "Mevsimsel şekli ne kadar yakaladık."],
    ["Sezon MAE / inhouse", "Sezon hatası / sezon hacmi", "Sezon içi planlama kalitesi."],
])
metric_table.columns = ["Metrik", "Formül", "Ne zaman kullanılır"]
add_table(slide, metric_table, Inches(0.65), Inches(1.35), Inches(11.9), Inches(2.6),
          col_widths=[0.18, 0.30, 0.52], size=13)
add_text(slide, Inches(0.65), Inches(4.3), Inches(11.9), Inches(2.6), [
    "Neden birden fazla metrik?",
    "- MAPE düşük hacimli aylarda anlamsızlaşır: Muğla Ocak ayı ~30 gece, 10 gecelik sapma %33 hata.",
    "- Aynı birimde MAPE %500, WMAPE %27 olabilir (örn. Manavgat 2026 Ridge backtesti).",
    "- Bu yüzden karar metriği olarak WMAPE ve toplam hata yüzdesi öne çıkarılıyor.",
], size=14)
notes(slide, "Bu slayt tartışmayı yönetir. Yüksek MAPE gördüğünde panik yok — hacme bak.")

# ======================================================================================
# ANTALYA
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = ANT
bg.line.fill.background(); bg.shadow.inherit = False
box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5))
frame = box.text_frame
frame.text = "BÖLÜM 1 · ANTALYA"
frame.paragraphs[0].runs[0].font.size = Pt(44)
frame.paragraphs[0].runs[0].font.bold = True
frame.paragraphs[0].runs[0].font.color.rgb = WHITE
para = frame.add_paragraph()
para.text = "Şehir · 12 ilçe · Top 20 otel"
para.runs[0].font.size = Pt(18)
para.runs[0].font.color.rgb = RGBColor(0xFF, 0xE8, 0xD9)
notes(slide, "Antalya bölümüne geçiş.")

slide_image(
    "Antalya · Aylık Inhouse Serisi",
    "2023-01 … 2026-07, occupied nights",
    "01_monthly_time_series.png",
    [
        "Gözlemler",
        "- Çok güçlü ve düzenli yıllık mevsimsellik.",
        "- Temmuz-Ağustos zirve, Aralık-Şubat dip.",
        "- 2023 → 2025 arasında net hacim büyümesi.",
        "",
        "Modelleme açısından",
        "- Uzun, düzenli ve yüksek hacimli seri.",
        "- Mevsimsel farklanma gerekiyor (durağan değil).",
        "- Bu yapı SARIMAX için elverişli.",
    ],
    "Antalya serisi modellemek için ideal: düzenli, yüksek hacimli, net mevsimsellik.",
    color=ANT,
)

slide_image(
    "Antalya · Yıl Bazlı Mevsimsel Eğriler",
    "Aylar x ekseninde, yıllar üst üste — takvim kaymasının görüldüğü grafik",
    "01_seasonal_curves.png",
    [
        "KRİTİK GÖZLEM",
        "- 2023, 2024 ve 2025 eğrileri neredeyse aynı şekle sahip.",
        "- 2026 eğrisi Mayıs'ta beklenmedik bir sıçrama yapıyor,",
        "  ardından Haziran'da geriliyor.",
        "",
        "Sebep",
        "- Kurban Bayramı 2026'da Mayıs sonuna kaydı.",
        "- Talep Haziran'dan Mayıs'a öne çekildi.",
        "",
        "Sonuç",
        "- Sabit aylık mevsimsellik öğrenen model burada şaşar.",
        "- Düzeltme gerekiyor → sonraki slayt.",
    ],
    "Sunumun en önemli görselı. 2026 eğrisinin diğerlerinden ayrıştığı noktayı ekranda göster.",
    color=ANT,
)

slide_full_image(
    "Antalya · Trend ve Mevsimsellik Ayrıştırması",
    "Additive decomposition (period = 12): gözlem, trend, mevsimsellik, artık",
    "01_trend_seasonality.png",
    "Trend bileşeni istikrarlı yükseliyor, mevsimsel bileşen çok düzenli. Artıklar küçük — "
    "yani seri büyük ölçüde trend + mevsimsellik ile açıklanıyor.",
    color=ANT,
)

slide_bullets(
    "Antalya · Takvim Etkisi ve Uygulanan Düzeltme",
    "Modelin en büyük hata kaynağı ve nasıl telafi edildi",
    [
        "Sorun: hicri takvim her yıl ~11 gün geriye kayar; sabit aylık mevsimsellik bunu göremez.",
        "- 2025 Kurban: Haziran başı → 2026 Kurban: Mayıs sonu → 2027 Kurban: 16-19 Mayıs",
        "- 2025 Ramazan: 30 Mart → 2026 Ramazan: 20 Mart → 2027 Ramazan: 9-11 Mart",
        "",
        "Uygulanan dört düzeltme",
        "- Mayıs uplift: baseline'ın Mayıs'ta kaçırdığı hacim kadar ekleme",
        "- Haziran transfer faktörü = 0,648 → Haziran tahmini %35 aşağı çekiliyor",
        "- Mart uplift ve Nisan transfer: Ramazan kaymasının aynı mantıkla telafisi",
        "",
        "Etki (2026 Ocak-Temmuz testi)",
        "- SARIMAX baseline WMAPE %36,95 → düzeltmeli %17,62  (yarı yarıya azalma)",
        "- MAE 3.968 → 1.892 gece   |   R² 0,850 → 0,975",
        "- Prophet'te de aynı yönde iyileşme: %42,18 → %23,16",
    ],
    "Sunumun kalbi. Hatanın rastgele değil, açıklanabilir ve düzeltilebilir bir takvim etkisi "
    "olduğunu anlat. Rakamları ezberle: 36,95 → 17,62.",
    size=14.5,
    color=ANT,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Antalya · 2026 Backtest Sonuçları", "Ocak-Temmuz 2026, dört model karşılaştırması", color=ANT)
path = img("03_test_predictions_comparison.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(7.6), Inches(3.7))
antalya_models = pd.DataFrame([
    ["SARIMAX + takvim düzeltmesi", "1.892", "17,62", "0,975"],
    ["Prophet + takvim düzeltmesi", "2.486", "23,16", "0,938"],
    ["SARIMAX baseline", "3.968", "36,95", "0,850"],
    ["Prophet baseline", "4.529", "42,18", "0,792"],
])
antalya_models.columns = ["Model", "MAE (gece)", "WMAPE %", "R²"]
add_table(slide, antalya_models, Inches(0.5), Inches(5.15), Inches(7.3), Inches(1.6), size=12)
add_text(slide, Inches(8.2), Inches(1.35), Inches(4.8), Inches(5.7), [
    "Nasıl okunur",
    "- Siyah kalın çizgi: gerçek 2026 değerleri.",
    "- Turuncu: takvim düzeltmeli SARIMAX.",
    "- Mavi: takvim düzeltmeli Prophet.",
    "",
    "Sonuç",
    "- SARIMAX + düzeltme açık ara en iyi (%17,6 WMAPE).",
    "- R² 0,975 → mevsimsel şekil neredeyse tam yakalanmış.",
    "- Prophet düzeltmeyle bile SARIMAX'ın gerisinde.",
    "",
    "Not",
    "- Düzeltme katsayıları 2026 gerçeğinden türetildi.",
    "- Metodoloji uyarısı için 18. slayta bakın.",
], size=12.5)
notes(slide, "Bu tabloyu ekranda göster. Ama %17,6'nın nasıl elde edildiğini 18. slaytta dürüstçe "
             "açıklayacağını da belirt.")

slide_image(
    "Antalya · Aylık Hata Dağılımı",
    "Gerçek − tahmin, düzeltme uygulanmış modeller",
    "03_forecast_errors.png",
    [
        "Gözlemler",
        "- Hatalar sıfır çizgisinin etrafında dağılıyor.",
        "- Sistematik bir yön (hep fazla/hep eksik) yok.",
        "",
        "Bu neden iyi bir işaret?",
        "- Sistematik sapma olsaydı model yanlı olurdu.",
        "- Dağınık hata, kalan kısmın büyük ölçüde",
        "  öngörülemez gürültü olduğunu gösterir.",
        "",
        "Dikkat",
        "- Zirve aylarda mutlak hata doğal olarak büyür;",
        "  yüzdesel olarak ise küçüktür.",
    ],
    "Hata grafiği modelin yanlı olmadığını gösterir; bu güven verici bir kanıt.",
    color=ANT,
)

# --- Antalya ilçe
slide_full_image(
    "Antalya · İlçe Bazlı Sezon / Sezon Dışı Hacim",
    "Aktif ilçeler, 2023-2025 · Sezon: Mart-Ekim, Sezon dışı: Kasım-Şubat",
    "24_cell24_img1.png",
    "Hacim dağılımını göster: Serik, Manavgat, Kemer ve Alanya Antalya talebinin büyük kısmını "
    "taşıyor. Aksu, Döşemealtı, Kepez neredeyse sıfır.",
    color=ANT,
)

slide_full_image(
    "Antalya · İlçe Bazlı 2025 Backtest",
    "Eğitim 2023-2024 → Test 2025 · her ilçe için seçilen model ve metrikler",
    "26_filtered_active_district_2025_actual_vs_model_metrics.png",
    "Panellerde her ilçenin gerçek ve tahmin eğrisi, kutuda MAPE/WMAPE/R² ve sezon-sezon dışı "
    "hata oranı var. Kemer ve Alanya'nın uyumuna dikkat çek.",
    color=ANT,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Antalya · İlçe Tahmin Kalitesi Sıralaması", "2025 backtest, sezon hata oranına göre", color=ANT)
path = img("26_filtered_active_district_2025_error_ratio_rankings.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(7.4), Inches(5.9))
district_table = pd.DataFrame([
    ["Kemer", "SARIMAX", "22,3", "8,4", "0,95"],
    ["Kumluca", "Prophet", "32,7", "8,8", "0,88"],
    ["Serik", "SARIMAX", "31,5", "9,1", "0,88"],
    ["Muratpaşa", "SARIMAX", "30,4", "14,0", "0,70"],
    ["Alanya", "SARIMAX", "27,1", "22,5", "0,91"],
    ["Manavgat", "Prophet", "27,4", "23,2", "0,86"],
    ["Kaş", "Prophet", "38,4", "35,3", "0,79"],
])
district_table.columns = ["İlçe", "Model", "WMAPE %", "Sezon hata %", "R²"]
add_table(slide, district_table, Inches(8.0), Inches(1.45), Inches(4.95), Inches(3.1),
          col_widths=[0.24, 0.24, 0.20, 0.20, 0.12], size=11.5)
add_text(slide, Inches(8.0), Inches(4.8), Inches(4.95), Inches(2.3), [
    "Öne çıkanlar",
    "- Kemer sezon hacminin %8'i kadar hata ile en iyi.",
    "- Model dağılımı karışık: 4 SARIMAX, 3 Prophet.",
    "- Yani tek model her ilçeye uymuyor.",
], size=12)
notes(slide, "Model dağılımının karışık olması bilinçli tasarımın sonucu — her ilçe kendi en iyi "
             "modelini kullanıyor.")

slide_full_image(
    "Antalya · Seçili 7 İlçe, 2025 Aylık Tahmin",
    "Alanya · Kaş · Kemer · Kumluca · Manavgat · Muratpaşa · Serik",
    "27_selected_7_districts_2025_actual_vs_model_metrics.png",
    "Karar için en anlamlı 7 ilçe. Her panelde seçilen model, MAPE, WMAPE, R² ve sezon/sezon dışı "
    "hata oranları yazılı.",
    color=ANT,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Antalya · İlçe 2026 Mart-Temmuz Backtesti", "Ridge (ay etkileri + yıl trendi) · eğitim 2023-2025", color=ANT)
path = img("29_district_monthly_mar_jul_2026_actual_vs_prediction.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(8.0), Inches(5.9))
ridge_table = pd.DataFrame([
    ["Kumluca", "26,2", "191,5", "0,93"],
    ["Muratpaşa", "51,0", "114,2", "0,46"],
    ["Alanya", "53,0", "224,0", "0,66"],
    ["Serik", "60,5", "164,8", "0,54"],
    ["Kaş", "61,6", "251,7", "0,65"],
    ["Manavgat", "64,1", "499,5", "0,54"],
    ["Kemer", "66,6", "575,3", "0,66"],
])
ridge_table.columns = ["İlçe", "WMAPE %", "MAPE %", "R²"]
add_table(slide, ridge_table, Inches(8.55), Inches(1.45), Inches(4.4), Inches(3.1),
          col_widths=[0.34, 0.22, 0.24, 0.20], size=11.5)
add_text(slide, Inches(8.55), Inches(4.8), Inches(4.4), Inches(2.3), [
    "Neden 2025'ten kötü?",
    "- Test dönemi tam da bayramın kaydığı aylar (Mar-Tem).",
    "- Ridge modeline takvim düzeltmesi uygulanmadı.",
    "- MAPE'lerin %500'e çıkması düşük hacimli",
    "  Mart-Nisan aylarından kaynaklanıyor.",
], size=12)
notes(slide, "Bu slayt takvim düzeltmesinin değerini dolaylı kanıtlıyor: düzeltme uygulanmayan "
             "Ridge modeli aynı dönemde çok daha kötü.")

# --- Antalya otel
slide_full_image(
    "Antalya · Top 20 Otel, Aylık Inhouse Isı Haritası",
    "2023-01 … 2026-07, occupied nights",
    "14_top20_hotels_monthly_inhouse_heatmap.png",
    "Otellerin sezon yoğunluğunu ve birbirinden farklı doluluk profillerini gösterir. "
    "Koyu bantlar Temmuz-Ağustos yoğunluğu.",
    color=ANT,
)

slide_full_image(
    "Antalya · Top 20 Otel, Yıl Bazlı Ocak-Temmuz Hacmi",
    "2023 · 2024 · 2025 · 2026 karşılaştırması",
    "14_top20_hotels_january_july_by_year.png",
    "Hangi otelin büyüdüğü, hangisinin daraldığı burada görülür. Otel bazlı ticari görüşmelerde "
    "kullanılabilecek bir görsel.",
    color=ANT,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Antalya · Otel Bazlı 2026 Backtest", "Filtre: 2023-2025 Ocak-Temmuz inhouse ≥ 200 gece", color=ANT)
path = img("16_top20_2026_actual_vs_best_model_with_metrics.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(8.1), Inches(5.9))
hotel_table = pd.DataFrame([
    ["Sueno Hotels Deluxe Belek", "Prophet", "31,9", "0,89"],
    ["Robinson Çamyuva", "SARIMAX", "32,5", "0,89"],
    ["Belek Beach Resort", "SARIMAX", "36,3", "0,82"],
    ["Crystal Admiral Aqua", "Prophet", "44,8", "0,71"],
    ["Swandor Topkapı Palace", "Prophet", "48,1", "0,45"],
    ["Crystal Waterworld", "Prophet", "59,2", "0,62"],
    ["Champion Holiday Village", "SARIMAX", "70,1", "0,38"],
])
hotel_table.columns = ["Otel", "Model", "WMAPE %", "R²"]
add_table(slide, hotel_table, Inches(8.65), Inches(1.45), Inches(4.3), Inches(3.0),
          col_widths=[0.44, 0.24, 0.18, 0.14], size=10.5)
add_text(slide, Inches(8.65), Inches(4.7), Inches(4.3), Inches(2.4), [
    "Yorum",
    "- En iyi oteller %32-36 bandında.",
    "- Otel seviyesinde hata şehirden yüksek —",
    "  bu beklenen davranış (agregasyon gürültüyü söndürür).",
    "- %60 üzeri otellerde model önerilmez.",
], size=11.5)
notes(slide, "Otel seviyesinde tek tek gezinme. 'Her otel için ayrı model seçtik' mesajını ver.")

slide_full_image(
    "Antalya · Top 30 Otel, Sezon vs Sezon Dışı Hacim",
    "2023-2025 · Sezon: Mart-Ekim (sol), Sezon dışı: Kasım-Şubat (sağ)",
    "17_top30_season_vs_offseason_heatmaps.png",
    "Antalya otellerinin sezon dışında da (az da olsa) satışı var. Bu, Muğla ile en büyük farklardan biri.",
    color=ANT,
)

# ======================================================================================
# MUĞLA
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = MUG
bg.line.fill.background(); bg.shadow.inherit = False
box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5))
frame = box.text_frame
frame.text = "BÖLÜM 2 · MUĞLA"
frame.paragraphs[0].runs[0].font.size = Pt(44)
frame.paragraphs[0].runs[0].font.bold = True
frame.paragraphs[0].runs[0].font.color.rgb = WHITE
para = frame.add_paragraph()
para.text = "Şehir · 10 ilçe · filtreli 12 otel"
para.runs[0].font.size = Pt(18)
para.runs[0].font.color.rgb = RGBColor(0xD5, 0xE6, 0xF7)
notes(slide, "Muğla bölümüne geçiş.")

slide_full_image(
    "Muğla · Yıl Bazlı Mevsimsel Eğriler ve Yıllık Toplam",
    "2022-2026 · sol: aylara göre yıl eğrileri, sağ: Ocak-Temmuz toplamı",
    "34_cell34_img1.png",
    "Muğla'nın mevsimselliği Antalya'dan daha keskin: sezon dışı neredeyse sıfır. "
    "Ocak-Temmuz toplamında 2026 belirgin sıçrama gösteriyor.",
    color=MUG,
)

slide_full_image(
    "Muğla · Sezon ve Sezon Dışı Inhouse (2022-2026)",
    "Sol: yıl bazlı sezon/sezon dışı · Sağ: toplam",
    "32_cell32_img1.png",
    "Sezon dışı hacim toplamın çok küçük bir kısmı. Bu, yüzdesel metrikleri bozan yapısal bir özellik.",
    color=MUG,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Muğla · 2026 Backtest Sonuçları", "Ocak-Temmuz 2026, dört model + segment metrikleri", color=MUG)
path = img("38_cell38_img1.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(8.0), Inches(5.9))
mugla_models = pd.DataFrame([
    ["SARIMAX + takvim düzeltmesi", "790", "19,53", "0,916"],
    ["Prophet + takvim düzeltmesi", "976", "24,12", "0,928"],
    ["SARIMAX baseline", "1.143", "28,24", "0,888"],
    ["Prophet baseline", "1.176", "29,05", "0,912"],
])
mugla_models.columns = ["Model", "MAE", "WMAPE %", "R²"]
add_table(slide, mugla_models, Inches(8.55), Inches(1.45), Inches(4.4), Inches(1.9),
          col_widths=[0.44, 0.18, 0.22, 0.16], size=10.5)
add_text(slide, Inches(8.55), Inches(3.6), Inches(4.4), Inches(3.5), [
    "Segment bazlı (SARIMAX düzeltmeli)",
    "- Sezon (28.237 gece): WMAE %19,5",
    "- Sezon dışı (91 gece): WMAE %32,8",
    "",
    "Takvim düzeltme katsayıları",
    "- Mayıs uplift: +2.467 gece",
    "- Haziran transfer: 1,000 (Muğla'da düşüş yok)",
    "",
    "Yorum",
    "- Muğla'da bayram etkisi Antalya'dakinin tersine",
    "  Haziran'ı düşürmemiş, Temmuz'u yükseltmiş.",
], size=12)
notes(slide, "Muğla'da Haziran transfer faktörü 1,000 çıktı — yani talep Haziran'dan çekilmemiş. "
             "İki şehrin bayrama tepkisi farklı; bu ilginç bir bulgu.")

slide_full_image(
    "Muğla · İlçe Bazlı Hacim ve Isı Haritası",
    "2022-2026 Temmuz, ilk 10 ilçe",
    "35_cell35_img1.png",
    "Bodrum açık ara lider; Marmaris ve Fethiye onu takip ediyor. Muğla talebi Antalya'ya göre "
    "çok daha konsantre.",
    color=MUG,
)

slide_full_image(
    "Muğla · İlçe Bazlı Aylık Trend (55 Ay)",
    "2022-01 … 2026-07, en yüksek hacimli 8 ilçe",
    "36_cell36_img1.png",
    "Bodrum'un zirvelerinin diğerlerini nasıl gölgede bıraktığını göster. Marmaris'in 2026'daki "
    "düşüşü de bu grafikte görülüyor.",
    color=MUG,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Muğla · İlçe Bazlı 2026 Backtest", "Her ilçe için SARIMAX/Prophet + takvim düzeltmesi", color=MUG)
path = img("40_cell40_img1.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(7.9), Inches(5.9))
mugla_district = pd.DataFrame([
    ["Datça", "SARIMAX_adj", "15,2", "4,2", "0,99"],
    ["Fethiye", "Prophet_adj", "16,3", "4,6", "0,97"],
    ["Ortaca", "SARIMAX_adj", "23,7", "0,8", "0,94"],
    ["Bodrum", "SARIMAX_adj", "27,2", "21,0", "0,84"],
    ["Ula", "Prophet_adj", "28,2", "7,2", "0,89"],
    ["Milas", "SARIMAX_adj", "85,2", "32,6", "0,04"],
    ["Marmaris", "SARIMAX_adj", "201,1", "192,3", "-10,8"],
])
mugla_district.columns = ["İlçe", "Model", "WMAPE %", "Sezon top. hata %", "R²"]
add_table(slide, mugla_district, Inches(8.45), Inches(1.45), Inches(4.5), Inches(3.1),
          col_widths=[0.22, 0.26, 0.18, 0.22, 0.12], size=10.5)
add_text(slide, Inches(8.45), Inches(4.8), Inches(4.5), Inches(2.3), [
    "DİKKAT · Marmaris",
    "- WMAPE %201, R² −10,8 → model tamamen şaştı.",
    "- 2026'da Marmaris talebi sert düştü;",
    "  geçmişte örneği olmayan yapısal kırılma.",
    "- Bu ilçe için model çıktısı kullanılmamalı.",
], size=11.5)
notes(slide, "Marmaris'i dürüstçe anlat — modelin sınırını gösteren iyi bir örnek. "
             "Datça ve Fethiye ise çok iyi çalışıyor.")

slide_full_image(
    "Muğla · İlçe Bazlı Sezonsal Panel Karşılaştırması",
    "Mart-Temmuz 2026 · her ilçe için gerçek, tahmin ve metrik kutusu",
    "41_cell41_img1.png",
    "Her ilçenin ayrı paneli. Ortaca'nın toplam hatası %0,8 ile en iyi; Marmaris'in sapması "
    "gözle görülür.",
    color=MUG,
)

slide = prs.slides.add_slide(BLANK)
add_band(slide, "Muğla · Otel Bazlı 2026 Backtest", "Filtre: 2022-2025 her sezon ≥ 100 gece → 30 otelden 12'si", color=MUG)
path = img("49_mugla_12_2026_jan_jul_sarimax_may_june_backtest_metrics.png")
if path:
    add_picture_fit(slide, path, Inches(0.35), Inches(1.2), Inches(8.1), Inches(5.9))
mugla_hotel = pd.DataFrame([
    ["Divan Bodrum", "18,8", "15,3"],
    ["Costa Farilya", "23,8", "16,9"],
    ["Golden Age Yalıkavak", "35,6", "18,4"],
    ["Manaspark Deluxe", "37,1", "29,7"],
    ["Xanadu Island", "42,5", "42,4"],
    ["Robinson Sarıgerme", "50,6", "2,7"],
    ["Hillside Beach Club", "51,2", "6,7"],
])
mugla_hotel.columns = ["Otel", "WMAPE %", "Sezon top. hata %"]
add_table(slide, mugla_hotel, Inches(8.65), Inches(1.45), Inches(4.3), Inches(3.0),
          col_widths=[0.46, 0.24, 0.30], size=10.5)
add_text(slide, Inches(8.65), Inches(4.7), Inches(4.3), Inches(2.4), [
    "Dikkat çekici",
    "- Robinson Sarıgerme WMAPE %50 ama sezon",
    "  toplam hatası sadece %2,7.",
    "- Yani aylık dağılım şaşmış, toplam tutmuş.",
    "- Kontenjan planlaması için toplam hata daha önemli.",
], size=11.5)
notes(slide, "Robinson Sarıgerme örneği metrik seçiminin neden önemli olduğunu somutlaştırır — "
             "5. slayta geri bağla.")

slide_full_image(
    "Muğla · Filtreli 12 Otel, Sezon ve Sezon Dışı Trend",
    "2022-2026 · turuncu: sezon (Mar-Eki), mavi: sezon dışı (Kas-Şub)",
    "50_cell50_img1.png",
    "Otellerin sezon dışı çizgisinin neredeyse sıfırda yatay olduğunu göster — Muğla'nın yapısal özelliği.",
    color=MUG,
)

# ======================================================================================
# KARŞILAŞTIRMA
# ======================================================================================
slide = prs.slides.add_slide(BLANK)
add_band(slide, "Antalya vs Muğla · Seviye Seviye Karşılaştırma", "2026 backtest, en iyi modeller")
compare = pd.DataFrame([
    ["Şehir (2026 Oca-Tem)", "17,6", "19,5", "Antalya biraz önde"],
    ["Şehir · düzeltmesiz", "37,0", "28,2", "Muğla baseline'da daha iyi"],
    ["Düzeltme kazancı (puan)", "19,3", "8,7", "Antalya'da etki çok daha büyük"],
    ["En iyi ilçe", "Kemer %22,3", "Datça %15,2", "Muğla'nın en iyisi daha iyi"],
    ["En kötü ilçe", "Kaş %38,4", "Marmaris %201", "Muğla'da uç sapma var"],
    ["En iyi otel", "Sueno Belek %31,9", "Divan Bodrum %18,8", "Muğla'nın en iyisi daha iyi"],
])
compare.columns = ["Seviye / Metrik (WMAPE %)", "Antalya", "Muğla", "Yorum"]
add_table(slide, compare, Inches(0.6), Inches(1.35), Inches(12.1), Inches(2.9),
          col_widths=[0.28, 0.18, 0.18, 0.36], size=12.5)
add_text(slide, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.6), [
    "Kritik ayrım: ortalama mı, uç değerler mi?",
    "- Muğla'nın EN İYİ birimleri Antalya'nınkinden daha iyi (Datça %15, Divan Bodrum %19).",
    "- Ama Muğla'nın EN KÖTÜ birimi çok daha kötü (Marmaris %201, R² −10,8).",
    "- Yani Muğla'da varyans yüksek: bazı ilçelerde mükemmel, bazılarında kullanılamaz.",
    "- Antalya daha dengeli: en iyi %22, en kötü %38 — dar bir bant.",
    "İş kararı açısından: Antalya'da modele genel olarak güvenilebilir; Muğla'da birim birim seçim şart.",
], size=14)
notes(slide, "Bu slaytta 'hangisi daha iyi' sorusunun tek cevabı olmadığını anlat: ortalamaya "
             "bakarsan Antalya, en iyi birimlere bakarsan Muğla.")

slide_bullets(
    "Neden Bu Farklar Oluşuyor?",
    "İki şehrin talep karakteri ve modele etkisi",
    [
        "1) Hacim ve gürültü",
        "- Antalya rezervasyon hacmi Muğla'nın yaklaşık 2,5 katı → oransal gürültü daha düşük.",
        "- Yüksek hacim, mevsimsel yapının daha net öğrenilmesini sağlıyor.",
        "",
        "2) Talebin dağılımı",
        "- Antalya: talep 12 ilçeye yayılmış (Serik, Manavgat, Alanya, Kemer dengeli paylaşıyor).",
        "- Muğla: Bodrum tek başına baskın; Marmaris ve Fethiye ikinci grup, kalanlar çok küçük.",
        "- Konsantrasyon, tek bir ilçedeki kırılmanın toplamı bozmasına yol açıyor (Marmaris örneği).",
        "",
        "3) Sezon uzunluğu",
        "- Antalya'da sezon dışı da satış var (Ocak ~1.000 gece mertebesinde).",
        "- Muğla'da sezon dışı neredeyse sıfır (Ocak ~30 gece) → yüzdesel metrikler patlıyor.",
        "",
        "4) Bayrama tepki farkı",
        "- Antalya: Kurban Mayıs'a kayınca Haziran çöktü (transfer faktörü 0,648).",
        "- Muğla: Haziran düşmedi (faktör 1,000), etki Temmuz zirvesini büyüttü.",
        "- Aynı takvim olayı iki destinasyonda farklı davranış üretiyor → ayrı kalibrasyon gerekli.",
    ],
    "Sebep-sonuç slaytı. 'Model Antalya'da daha iyi' demek yerine 'Antalya'nın verisi modellemeye "
    "daha uygun' de. Fark modelden değil veriden geliyor.",
    size=14,
)

# ======================================================================================
# METODOLOJİ UYARILARI
# ======================================================================================
slide_bullets(
    "Metodolojik Uyarı 1 · Kalibrasyon Sızıntısı",
    "%17,6 rakamı nasıl okunmalı?",
    [
        "Takvim düzeltme katsayıları şu şekilde hesaplandı:",
        "- Mayıs uplift = 2026 Mayıs GERÇEĞİ − baseline tahmini",
        "- Haziran transfer = 2026 Haziran GERÇEĞİ / geçmiş Haziran ortalaması",
        "",
        "Yani düzeltme, test döneminin gerçek değerlerini kullanıyor.",
        "- Bu nedenle %17,62 kör (out-of-sample) bir tahmin başarısı DEĞİLDİR.",
        "- Doğru okunuşu: 'Bayram kaymasını bildiğimizi varsayarsak kalan hata %17,6'dır.'",
        "- Gerçek kör tahmin performansı baseline değerleridir: Antalya %37,0, Muğla %28,2.",
        "",
        "Bu bir hata mı? Hayır — ama etiketlenmesi şart.",
        "- Amaç, hatanın ne kadarının takvim kaynaklı olduğunu ölçmekti; bu ölçüm başarılı.",
        "- 2027 tahmininde aynı katsayılar ileriye taşınıyor; bu makul çünkü bayram tarihleri bilinir.",
        "",
        "Öneri: sunumda iki rakamı birlikte verin.",
        "- 'Düzeltmesiz %37 → takvim etkisini modellediğimizde %17,6'ya iniyor.'",
    ],
    "ÖNEMLİ: Bu slaydı atlama. Biri 'bu rakam nasıl bu kadar iyi?' diye sorarsa hazırlıklı ol. "
    "Kendin söylersen güven kazanırsın, sorulursa savunmaya geçersin.",
    size=14.5,
    color=WARN,
)

slide_bullets(
    "Metodolojik Uyarı 2 · Uzun Ufuk ve 2027 Tahmini",
    "SARIMAX ileri ufukta neden kullanılamadı?",
    [
        "Notebook çıktısındaki bulgu:",
        "- SARIMAX ile üretilen 2027 Haziran tahmini: −48.037 gece (negatif).",
        "- Occupied nights negatif olamaz → model uzun ufukta patlıyor.",
        "",
        "Kök neden",
        "- order=(1,2,1) ile ikinci dereceden farklanma + sabit trend terimi,",
        "  ileri ufukta doğrusal trendi kontrolsüz biçimde uzatıyor.",
        "",
        "Uygulanan çözüm",
        "- 2027 tahmini için Prophet baseline'a geçildi (pozitif kalıyor).",
        "- Takvim düzeltmeleri (Mayıs uplift, Haziran/Nisan transfer, Mart uplift) Prophet üzerine eklendi.",
        "- 2027 Mayıs için ek bir taban kuralı: tahmin 2026 Mayıs gerçeğinin altına düşemez.",
        "",
        "Kalan risk",
        "- Bu bir yama; modelin kendisi düzeltilmedi.",
        "- Kalıcı çözüm: trend terimini kaldırmak (trend='n') veya log dönüşümü uygulamak,",
        "    tahmin ufkunu 6 aya indirip her ay yeniden eğitmek (rolling refit).",
        "- 2027 çıktısı yayınlanmadan önce sezon dışı aylar geçmiş ortalama ile kıyaslanmalı.",
    ],
    "İkinci uyarı. Sorunu bulduğunu ve nasıl yamaladığını anlat, kalıcı çözümü de öner. "
    "Bu, olgun bir analiz izlenimi verir.",
    size=14,
    color=WARN,
)

# ======================================================================================
# KAPANIŞ
# ======================================================================================
slide_bullets(
    "Riskler ve Modelin Sınırları",
    "Dürüst değerlendirme",
    [
        "Veri sınırları",
        "- Antalya için 3,5 yıllık, Muğla için 4,5 yıllık geçmiş; mevsimsellik için ideal olan 6+ yıl değil.",
        "- 2026 verisi Temmuz'da bitiyor; yıl sonu davranışı gözlenemedi.",
        "- Muğla ilçe eşlemesi anahtar kelime tabanlı; resmi ilçe master'ı ile doğrulanmalı.",
        "",
        "Model sınırları",
        "- Fiyat, kampanya, kur, rakip kapasitesi gibi dışsal değişkenler modelde yok.",
        "- Yapısal kırılmalar öngörülemiyor (Marmaris 2026: WMAPE %201, R² −10,8).",
        "- SARIMAX uzun ufukta negatif üretiyor; 2027 için Prophet yaması kullanıldı.",
        "",
        "Kullanım sınırları",
        "- Düşük hacimli birimlerde yüzdesel metrikler yanıltıcı; mutlak hataya bakılmalı.",
        "- WMAPE %60 üzeri ilçe/otellerde model çıktısı karar için kullanılmamalı.",
        "- Takvim düzeltmesi, bayram etkisinin geçmiştekiyle benzer tekrarlanacağı varsayımına dayanıyor.",
    ],
    "Riskleri sen söyle. Özellikle Marmaris, dışsal değişken eksikliği ve negatif tahmin konusunu.",
    size=14,
)

slide_bullets(
    "Öneriler ve Sonraki Adımlar",
    "Kısa ve orta vadeli aksiyon planı",
    [
        "Hemen kullanılabilir",
        "- Antalya ve Muğla şehir seviyesi tahminleri sezon planlamasında referans alınabilir.",
        "- Antalya'da Kemer, Kumluca, Serik, Muratpaşa, Alanya, Manavgat ilçeleri (sezon hatası ≤ %23).",
        "- Muğla'da Datça, Fethiye, Ortaca, Ula ilçeleri (sezon toplam hatası ≤ %7).",
        "- Otel seviyesinde: Antalya'da ilk 3, Muğla'da ilk 4 otel.",
        "",
        "Kullanmadan önce düzelt",
        "- 2027 ileri tahmini: sezon dışı ay kontrolü yapılmadan yayınlanmamalı.",
        "- Marmaris ve düşük hacimli birimler için model yerine ticari ekip görüşü.",
        "",
        "Model iyileştirme yol haritası",
        "- Bayram tarihlerini gerçek takvimden dışsal değişken (exog) olarak modele ekle → sızıntı biter.",
        "- SARIMAX'ta trend terimini kaldır veya log dönüşümü uygula → negatif tahmin sorunu çözülür.",
        "- Tahmin ufkunu 6 aya indir, her ay yeni gerçekleşmeyle yeniden eğit (rolling refit).",
        "- Fiyat ve kampanya verisini modele bağla; talep esnekliğini yakala.",
        "- Otomatik akıl testi ekle: geçmiş ay ortalamasından 3 kat sapma → uyarı üret.",
    ],
    "Kapanış. Somut ve sahiplenilebilir aksiyonlar ver, soru-cevaba buradan geç.",
    size=14,
)

slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK
bg.line.fill.background(); bg.shadow.inherit = False
box = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.5))
frame = box.text_frame
frame.word_wrap = True
frame.text = "Teşekkürler · Sorular"
frame.paragraphs[0].runs[0].font.size = Pt(40)
frame.paragraphs[0].runs[0].font.bold = True
frame.paragraphs[0].runs[0].font.color.rgb = WHITE
for text in [
    "Analiz defteri: notebooks/SARIMA.ipynb",
    "Sunum görselleri: SUNUM1/gorseller/ (57 grafik)",
    "Sunum planı ve konuşma notları: SUNUM1/SUNUM_PLANI.md",
]:
    para = frame.add_paragraph()
    para.text = text
    para.space_before = Pt(9)
    para.runs[0].font.size = Pt(14.5)
    para.runs[0].font.color.rgb = RGBColor(0xAD, 0xB5, 0xBD)
notes(slide, "Soru-cevap. Sık sorular: '%17,6 nasıl?' → 18. slayt. 'Hangi şehir daha iyi?' → "
             "16. slayt. 'Marmaris ne oldu?' → 14. slayt.")

output = SUNUM_DIR / "Antalya_Mugla_Talep_Analizi_SUNUM1.pptx"
prs.save(output)
print(f"Slayt sayisi: {len(prs.slides._sldIdLst)}")
print(f"Kaydedildi: {output}")
